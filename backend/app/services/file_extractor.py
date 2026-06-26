import os


def extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_md(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_xlsx(file_path: str) -> str:
    import pandas as pd
    dfs = pd.read_excel(file_path, sheet_name=None)
    parts = []
    for name, df in dfs.items():
        parts.append(f"=== Sheet: {name} ===\n{df.to_csv(index=False)}")
    return "\n\n".join(parts)


def extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return "\n".join(texts)


def extract_pdf(file_path: str) -> str:
    import pdfplumber
    pages = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
    except Exception as e:
        raise ValueError(f"PDF 解析失败（文件可能已损坏或不是有效的 PDF）: {e}")
    if not pages:
        raise ValueError("PDF 文件中未能提取到文本内容")
    return "\n\n".join(pages)


EXTRACTORS = {
    ".txt":  extract_txt,
    ".md":   extract_md,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".pdf":  extract_pdf,
}

ALLOWED_EXTENSIONS = set(EXTRACTORS.keys())


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"不支持的文件格式: {ext}")
    return extractor(file_path)
